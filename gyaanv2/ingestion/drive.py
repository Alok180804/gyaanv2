import io, re, tempfile
from pathlib import Path
import fitz
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from gyaanv2.models import LoadedDocument

SCOPES=["https://www.googleapis.com/auth/drive.readonly","https://www.googleapis.com/auth/documents.readonly","https://www.googleapis.com/auth/presentations.readonly","https://www.googleapis.com/auth/spreadsheets.readonly"]
FOLDER_MIME='application/vnd.google-apps.folder'


def parse_drive_link(url:str)->tuple[str,str]:
    patterns=[r'/folders/([a-zA-Z0-9_-]+)', r'/file/d/([a-zA-Z0-9_-]+)', r'/document/d/([a-zA-Z0-9_-]+)', r'/spreadsheets/d/([a-zA-Z0-9_-]+)', r'/presentation/d/([a-zA-Z0-9_-]+)', r'[?&]id=([a-zA-Z0-9_-]+)']
    for p in patterns:
        m=re.search(p,url)
        if m: return m.group(1), 'unknown'
    if re.fullmatch(r'[a-zA-Z0-9_-]+', url.strip()): return url.strip(), 'unknown'
    raise ValueError('Could not find a Google Drive id in the link.')

class GoogleDriveClient:
    def __init__(self, credentials_file:Path, token_file:Path):
        self.creds=self._auth(credentials_file, token_file)
        self.drive=build('drive','v3',credentials=self.creds)
        self.docs=build('docs','v1',credentials=self.creds)
        self.sheets=build('sheets','v4',credentials=self.creds)
        self.slides=build('slides','v1',credentials=self.creds)

    def _auth(self, credentials_file, token_file):
        creds=None
        if token_file.exists(): creds=Credentials.from_authorized_user_file(str(token_file), SCOPES)
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token: creds.refresh(Request())
            else:
                flow=InstalledAppFlow.from_client_secrets_file(str(credentials_file), SCOPES)
                creds=flow.run_local_server(port=0)
            token_file.write_text(creds.to_json())
        return creds

    def metadata(self, file_id):
        return self.drive.files().get(fileId=file_id, fields='id,name,mimeType,modifiedTime').execute()

    def iter_files(self, drive_id):
        meta=self.metadata(drive_id)
        if meta['mimeType']==FOLDER_MIME:
            token=None
            while True:
                res=self.drive.files().list(q=f"'{drive_id}' in parents and trashed=false", pageSize=1000, pageToken=token, fields='nextPageToken,files(id,name,mimeType,modifiedTime)').execute()
                for f in res.get('files',[]):
                    if f['mimeType']==FOLDER_MIME: yield from self.iter_files(f['id'])
                    else: yield f
                token=res.get('nextPageToken')
                if not token: break
        else:
            yield meta

    def load_file(self, source_id:int, file:dict):
        mime=file['mimeType']; fid=file['id']; name=file['name']; mod=file.get('modifiedTime')
        if mime=='application/vnd.google-apps.document':
            doc=self.docs.documents().get(documentId=fid).execute(); text=[]
            for el in doc.get('body',{}).get('content',[]):
                for item in el.get('paragraph',{}).get('elements',[]): text.append(item.get('textRun',{}).get('content',''))
            yield LoadedDocument(source_id,fid,name,mime,'google_doc',mod,''.join(text))
        elif mime=='application/vnd.google-apps.spreadsheet':
            book=self.sheets.spreadsheets().get(spreadsheetId=fid).execute()
            for sh in book.get('sheets',[]):
                sname=sh['properties']['title']; vals=self.sheets.spreadsheets().values().get(spreadsheetId=fid, range=sname).execute().get('values',[])
                text='\n'.join(' | '.join(map(str,row)) for row in vals)
                yield LoadedDocument(source_id,fid,name,mime,'google_sheet',mod,text,sheet_name=sname)
        elif mime=='application/vnd.google-apps.presentation':
            pres=self.slides.presentations().get(presentationId=fid).execute()
            for i,slide in enumerate(pres.get('slides',[]),1):
                parts=[]
                for el in slide.get('pageElements',[]):
                    for te in el.get('shape',{}).get('text',{}).get('textElements',[]): parts.append(te.get('textRun',{}).get('content',''))
                yield LoadedDocument(source_id,fid,name,mime,'google_slide',mod,''.join(parts),slide_number=i)
        else:
            data=self._download(fid)
            yield from self._load_binary(source_id,fid,name,mime,mod,data)

    def _download(self, fid):
        buf=io.BytesIO(); req=self.drive.files().get_media(fileId=fid); dl=MediaIoBaseDownload(buf, req); done=False
        while not done: _, done = dl.next_chunk()
        return buf.getvalue()

    def _load_binary(self, source_id,fid,name,mime,mod,data):
        suffix=Path(name).suffix.lower()
        if mime=='application/pdf' or suffix=='.pdf':
            pdf=fitz.open(stream=data, filetype='pdf')
            for i,page in enumerate(pdf,1):
                yield LoadedDocument(source_id,fid,name,mime,'pdf',mod,page.get_text(),page_number=i)
            pdf.close()
        elif suffix in ['.txt','.md','.csv']:
            yield LoadedDocument(source_id,fid,name,mime,suffix.lstrip('.'),mod,data.decode('utf-8', errors='ignore'))
        elif suffix=='.docx':
            from docx import Document as DocxDocument
            with tempfile.NamedTemporaryFile(suffix=suffix) as f:
                f.write(data); f.flush(); doc=DocxDocument(f.name)
                yield LoadedDocument(source_id,fid,name,mime,'docx',mod,'\n'.join(p.text for p in doc.paragraphs))
        elif suffix=='.xlsx':
            from openpyxl import load_workbook
            with tempfile.NamedTemporaryFile(suffix=suffix) as f:
                f.write(data); f.flush(); wb=load_workbook(f.name, read_only=True, data_only=True)
                for ws in wb.worksheets:
                    text='\n'.join(' | '.join('' if v is None else str(v) for v in row) for row in ws.iter_rows(values_only=True))
                    yield LoadedDocument(source_id,fid,name,mime,'xlsx',mod,text,sheet_name=ws.title)
        elif suffix=='.pptx':
            from pptx import Presentation
            with tempfile.NamedTemporaryFile(suffix=suffix) as f:
                f.write(data); f.flush(); prs=Presentation(f.name)
                for i,sl in enumerate(prs.slides,1):
                    text='\n'.join(getattr(shape,'text','') for shape in sl.shapes if hasattr(shape,'text'))
                    yield LoadedDocument(source_id,fid,name,mime,'pptx',mod,text,slide_number=i)
